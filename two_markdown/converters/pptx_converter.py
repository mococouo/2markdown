from __future__ import annotations

from pathlib import Path

from two_markdown.context import ConversionContext
from two_markdown.errors import MissingDependencyError
from two_markdown.models import ConvertedDocument
from two_markdown.utils import markdown_table, normalize_extracted_text


class PptxConverter:
    name = "pptx"
    extensions = {".pptx"}

    def convert(self, path: Path, context: ConversionContext) -> ConvertedDocument:
        try:
            from pptx import Presentation
        except Exception as exc:
            raise MissingDependencyError("python-pptx", "PPTX conversion") from exc

        presentation = Presentation(path)
        parts = [f"# {path.stem}"]
        warnings: list[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            parts.append(f"## Slide {index}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = _text_frame_to_markdown(shape.text_frame)
                    if text.strip():
                        parts.append(normalize_extracted_text(text))
                if getattr(shape, "has_table", False):
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    if rows:
                        parts.append(markdown_table(rows))
                if getattr(shape, "has_chart", False):
                    chart_md = _chart_to_markdown(shape.chart, warnings)
                    if chart_md:
                        parts.append(chart_md)
                try:
                    image = shape.image
                except Exception:
                    image = None
                if image is not None:
                    extension = image.ext or "png"
                    filename = f"slide-{index}-image.{extension}"
                    parts.append(context.add_attachment(filename, image.blob, alt_text=f"slide {index} image"))

            notes_text = _notes_text(slide)
            if notes_text:
                parts.append(f"### Notes\n\n{normalize_extracted_text(notes_text)}")

        if context.attachments:
            warnings.append("Images were extracted from slides.")

        metadata = _core_properties(presentation) if context.options.include_metadata else {}

        return ConvertedDocument(
            title=metadata.get("title") or path.stem,
            markdown="\n\n".join(parts),
            converter=self.name,
            attachments=context.attachments,
            warnings=warnings,
            metadata=metadata,
        )


def _text_frame_to_markdown(text_frame: object) -> str:
    lines: list[str] = []
    for paragraph in text_frame.paragraphs:  # type: ignore[attr-defined]
        chunks: list[str] = []
        for run in paragraph.runs:  # type: ignore[attr-defined]
            text = run.text or ""  # type: ignore[attr-defined]
            if not text:
                continue
            font = getattr(run, "font", None)
            bold = getattr(font, "bold", None) if font else None
            italic = getattr(font, "italic", None) if font else None
            underline = getattr(font, "underline", None) if font else None
            wrapped = text
            if bold:
                wrapped = f"**{wrapped}**"
            if italic:
                wrapped = f"*{wrapped}*"
            if underline:
                wrapped = f"<u>{wrapped}</u>"
            hyperlink = getattr(run, "hyperlink", None)
            address = getattr(hyperlink, "address", None) if hyperlink else None
            if address:
                wrapped = f"[{wrapped}]({address})"
            chunks.append(wrapped)
        line = "".join(chunks).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def _chart_to_markdown(chart: object, warnings: list[str]) -> str:
    try:
        chart_type = getattr(chart, "chart_type", None)
        parts = [f"**Chart** ({chart_type})"]
        for plot in chart.plots:  # type: ignore[attr-defined]
            categories = list(plot.categories or [])
            series_list = list(plot.series)
            header = ["Category"] + [_series_name(series) for series in series_list]
            rows: list[list[object]] = [header]
            for cat_index, category in enumerate(categories):
                row: list[object] = [category]
                for series in series_list:
                    values = list(series.values or [])
                    row.append(values[cat_index] if cat_index < len(values) else "")
                rows.append(row)
            parts.append(markdown_table(rows))
        return "\n\n".join(parts)
    except Exception as exc:
        warnings.append(f"Could not render chart: {exc}")
        return ""


def _series_name(series: object) -> str:
    name = getattr(series, "name", None)
    if name is None:
        return "Series"
    if hasattr(name, "ishtml"):
        try:
            name = name.text
        except Exception:
            name = str(name)
    return str(name).strip() or "Series"


def _notes_text(slide: object) -> str:
    try:
        notes_slide = slide.notes_slide  # type: ignore[attr-defined]
        text_frame = notes_slide.notes_text_frame  # type: ignore[attr-defined]
    except Exception:
        return ""
    lines = [paragraph.text.strip() for paragraph in text_frame.paragraphs if paragraph.text.strip()]  # type: ignore[attr-defined]
    return "\n".join(lines)


def _core_properties(presentation: object) -> dict[str, str]:
    metadata: dict[str, str] = {}
    try:
        props = presentation.core_properties  # type: ignore[attr-defined]
    except Exception:
        return metadata
    pairs = {
        "author": getattr(props, "author", None),
        "subject": getattr(props, "subject", None),
        "keywords": getattr(props, "keywords", None),
        "category": getattr(props, "category", None),
        "comments": getattr(props, "comments", None),
    }
    for key, value in pairs.items():
        if value and str(value).strip():
            metadata[key] = str(value).strip()
    title = getattr(props, "title", None)
    if title and str(title).strip():
        metadata["title"] = str(title).strip()
    created = getattr(props, "created", None)
    if created:
        metadata["created"] = _iso_dt(created)
    modified = getattr(props, "modified", None)
    if modified:
        metadata["modified"] = _iso_dt(modified)
    return metadata


def _iso_dt(value: object) -> str:
    iso = getattr(value, "isoformat", None)
    if callable(iso):
        try:
            return iso()
        except Exception:
            pass
    return str(value)
